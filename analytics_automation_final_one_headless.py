from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.chrome.options import Options
import os
import time
import re
from bs4 import BeautifulSoup
import tkinter as tk
from tkinter import ttk
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import sys
import shutil 
import csv

start_time = time.time()
 
# Folder functions
def folder_creation(user_info_dictionary):
    folder_path = os.path.join(os.path.expanduser("~"), "Desktop")
    folder_name = user_info_dictionary['event_name']

    folder_path = os.path.join(folder_path, folder_name)
    os.makedirs(folder_path)
    print(f'folder: {folder_name} created at{folder_path}')
def erase_folder(user_info_dictionary):
    folder_path = os.path.join(os.path.expanduser("~"), "Desktop")
    folder_name = user_info_dictionary['event_name']
    folder_path = os.path.join(folder_path, folder_name)
    shutil.rmtree(folder_path)

# ChromeDriver functions
def initialize_driver(download_folder):
    chromedriver_path = os.path.expanduser("~/Downloads/chromedriver")
    chrome_options = Options()
    chrome_options.add_argument('--window-size=1200,800')  # Adjust the window size
    chrome_options.add_argument('--headless')
    chrome_options.add_argument('--disable-gpu')  # necessary when running in headless mode
    chrome_options.add_argument('--user-agent=Your_User_Agent_String') # to simulate real user usage
    chrome_options.add_argument('--enable-logging') #additional loggin options
    chrome_options.add_argument('--log-level=0') #additional loggin options
    chrome_options.add_argument(f"executable_path={chromedriver_path}")
    # Set Chrome options to specify the download folder
    chrome_prefs = {"download.default_directory": download_folder}
    chrome_options.add_experimental_option("prefs", chrome_prefs)
    return webdriver.Chrome(options=chrome_options)
def handle_cookie_popup(driver):
    try:
        # Locate the "Allow All" button on the cookie popup and click it
        cookie_popup_button = WebDriverWait(driver, 10).until(
            EC.presence_of_element_located((By.ID, "CybotCookiebotDialogBodyLevelButtonLevelOptinAllowAll"))
        )
        cookie_popup_button.click()
    except Exception as e:
        print("Cookie popup not found or couldn't be handled:", e)
def login(driver, user_info_dictionary):
    # Locate the email input field and fill in the email
    email_element = WebDriverWait(driver, 20).until(
        EC.presence_of_element_located((By.ID, "e-mail"))
    )
    email_element.send_keys(user_info_dictionary['e_mail'])

    # Locate the password input field and fill in the password
    password_element = WebDriverWait(driver, 20).until(
        EC.presence_of_element_located((By.ID, "password"))
    )
    password_element.send_keys(user_info_dictionary['password'])

    # Locate and click the "Sign in" button
    sign_in_button = WebDriverWait(driver, 20).until(
        EC.element_to_be_clickable((By.XPATH, "//div[text()='Sign in']"))
    )
    sign_in_button.click()

    # Add a sleep delay to hold the script open for observation (one second is 
    #necessary for the login to work)
    time.sleep(1)  # Sleep for 1 seconds
def navigate_to_page(driver, user_info_dictionary):
    driver.get(user_info_dictionary['page_url'])
    WebDriverWait(driver, 20).until(
        EC.presence_of_element_located((By.TAG_NAME, "body"))
    )
def powered_by_navigation(driver):
    try:
        # Locate the "Allow All" button on the cookie popup and click it
        cookie_popup_button = WebDriverWait(driver, 10).until(
            EC.presence_of_element_located((By.ID, "CybotCookiebotDialogBodyLevelButtonLevelOptinAllowAll"))
        )
        cookie_popup_button.click()
    except Exception as e:
        print("Cookie popup not found or couldn't be handled:", e)

    time.sleep(1)

    # Click on the "Settings" button
    settings_button = WebDriverWait(driver, 3).until(
        EC.element_to_be_clickable((By.ID, "settings"))
    )
    settings_button.click()
    time.sleep(1)
    settings_button = WebDriverWait(driver, 3).until(
        EC.element_to_be_clickable((By.ID, "settings"))
    )
    settings_button.click()

    # Click on Analytics
    analyics_button = WebDriverWait(driver, 3).until(
        EC.element_to_be_clickable((By.XPATH, "//li[@class='nav-item']/a[@class='nav-link' and text()='Analytics']"))
    )
    analyics_button.click()
def enter_dates_and_submit(driver, user_info_dictionary):
    # Locate and fill in the start date field
    start_date_element = WebDriverWait(driver, 20).until(
        EC.presence_of_element_located((By.NAME, "sel_start"))
    )
    start_date_element.clear()  # Clear any existing value
    start_date_element.send_keys(user_info_dictionary['start_date'])

    # Locate and fill in the end date field
    end_date_element = WebDriverWait(driver, 20).until(
        EC.presence_of_element_located((By.NAME, "sel_stop"))
    )
    end_date_element.clear()  # Clear any existing value
    end_date_element.send_keys(user_info_dictionary['end_date'])

    # Locate and click the "Bericht erstellen" button
    submit_button = WebDriverWait(driver, 20).until(
        EC.element_to_be_clickable((By.CSS_SELECTOR, 'input[value="Generate report"]')) # may only work when on english
    )
    submit_button.click()

    time.sleep(4)
def download_transaction_data(driver, start_date, end_date, download_timeout):

    # Locate and fill in the start date field
    start_date_element = WebDriverWait(driver, 20).until(
        EC.presence_of_element_located((By.ID, "start_date"))
    )
    start_date_element.clear()  # Clear any existing value
    start_date_element.send_keys(start_date)

    # Locate and fill in the end date field
    end_date_element = WebDriverWait(driver, 20).until(
        EC.presence_of_element_located((By.NAME, "stop_date"))
    )
    end_date_element.clear()  # Clear any existing value
    end_date_element.send_keys(end_date)

    # Locate and click the "Bericht erstellen" button
    submit_button = WebDriverWait(driver, 20).until(
        EC.element_to_be_clickable((By.ID, 'submit-btn'))
    )
    submit_button.click()

    # Adjust waiting time based on download timeout
    driver.set_page_load_timeout(download_timeout)

# Function for information gathering
def capture_information(driver, info_type, user_info_dictionary):
    # Specify the full path to save the text document on the desktop
    desktop_path = os.path.expanduser(f"~/Desktop/{user_info_dictionary['event_name']}")

    # Construct the file path based on the information type
    text_file_path = os.path.join(desktop_path, f"captured_information_{info_type}.txt")

    # Wait for the page to load completely
    WebDriverWait(driver, 20).until(
        EC.presence_of_element_located((By.TAG_NAME, "body"))
    )

    # Get the page source
    page_source = driver.page_source

    # Use BeautifulSoup to parse the HTML
    soup = BeautifulSoup(page_source, "html.parser")

    # Save the prettified HTML content to a text document
    with open(text_file_path, "w", encoding="utf-8") as text_file:
        text_file.write(soup.prettify())

    print(f"Prettified HTML content saved to: {text_file_path}")

# Number formating Functions
def format_number(num):
    return '{:,.0f}'.format(num).replace(',', '.')  # Formats the number with thousands separat
def format_money(num, type):
    if len(type) == 2:
        num_str_1 = '{:,.2f}'.format(num / 100)  # Divide by 100 and format with two decimal places
        if num_str_1.endswith('.00'):  # Check if the last two digits are '.00'
            return num_str_1[:-3].replace(',', '.')  # Remove '.00' and replace comma
        else:
            return num_str_1.replace(',', '.')
    else:
        num_str_2 = '{:,}'.format(num)
        return num_str_2.replace(',', '.')
def reformat_numbers(value):
    value = str(value).replace(' ', '').replace(',','.').replace('EUR', '')
    if len(value) >= 3:
        if value[-3] == '.' or value[-2] == '.':
            while value.count('.') > 1:
                index = value.find('.')
                value = value[:index] + value[index+1:]
            value = float(value)
        elif value.find('.') != -1:
            value = value.replace('.', '')
            value = int(value)
        else:
            value = float(value)
    else:
        try:
            value = float(value)
        except:
            value = 1
    return value
# Value calculation functions
def calculate_analytic_values(unique_views, views, transaction_total_count,
                              transaction_all_access_pass_count,
                              transaction_singlestream_count,
                              total_brut_revenue, total_net_revenue, rev_share):
    values = [unique_views,views, transaction_total_count,
              transaction_all_access_pass_count, transaction_singlestream_count, 
              total_brut_revenue, total_net_revenue, rev_share]
    count = 0
    for value in values:
        values[count] = reformat_numbers(value)
        count += 1
    # Unique views / views
    view_ratio = round((values[0]/values[1]), 2)
    view_transaction_ratio = round((values[0]/values[2]), 2)
    singe_stream_all_access_pass_ratio = round((values[3]/values[4]), 2)
    avarage_vat = round(((values[5]*(values[7]/100)-values[6])/values[6]), 2)
    return (view_ratio, view_transaction_ratio, 
            singe_stream_all_access_pass_ratio, avarage_vat)
def calculate_single_stream_purchase_count(total_transactions, 
                                           all_access_pass_transactions):
    total_transactions = reformat_numbers(total_transactions)
    all_access_pass_transactions = reformat_numbers(all_access_pass_transactions)
    single_stream_purchase_count = (total_transactions-
                                    all_access_pass_transactions)
    single_stream_purchase_count = format_number(single_stream_purchase_count)
    return single_stream_purchase_count
def calculate_all_bought_all_access_passes(all_access_data):
    count = 0
    total_all_access_passes = 0
    while len(all_access_data)-count != 0:
        total_all_access_passes += int(all_access_data[count]['Active subscribers'])
        count +=1
    return total_all_access_passes

# Accessing the file depending on the data category you want
def access_file(data_category, user_info_dictionary):
    # Specify the path to your text document on the desktop
    desktop_path = f'~/Desktop/{user_info_dictionary['event_name']}/captured_information_{data_category}.txt'

    # Expand the ~ character to the home directory
    desktop_path = os.path.expanduser(desktop_path)

    # Read the content of the text document
    try:
        with open(desktop_path, 'r', encoding='utf-8') as file:
            html_text = file.read()
    except FileNotFoundError:
        print(f"File not found: {desktop_path}")
        exit()
    return html_text

# Extract specific info functions
def extract_views(html_text):
    soup = BeautifulSoup(html_text, 'html.parser')

    total_views_element = soup.find('span', id='total_views')
    unique_views_element = soup.find('span', id='unique_viewers')

    if not total_views_element or not unique_views_element:
        print("Views data not found in the HTML text.")
        return None, None

    total_views_value = re.sub(r'\D', '', total_views_element.get_text())
    unique_views_value = re.sub(r'\D', '', unique_views_element.get_text())

    def remove_every_second_digit(number):
        return int(''.join(number[i] for i in range(len(number)) if i % 2 != 0))

    # Check if the extracted values are not empty
    if total_views_value and unique_views_value:
    # Convert the values to integers and transform
        total_views_int = remove_every_second_digit(total_views_value)
        total_views_int = format_number(total_views_int)
        unique_views_int = remove_every_second_digit(unique_views_value)
        unique_views_int = format_number(unique_views_int)
    
    """ print(f"The total views are: {total_views_int}")
        print(f"The unique views are: {unique_views_int}")
    else:
        print("Views data not found in the HTML text.") """
    return total_views_int, unique_views_int
def extract_ppv_info(html_text):
    # Create a BeautifulSoup object with the HTML text
    soup = BeautifulSoup(html_text, 'html.parser')

    # Extract total transactions value
    total_transactions_element = soup.find('span', id='total_transactions')
    total_transactions_value = re.sub(r'\D', '', total_transactions_element.get_text())

    # Extract total transactions odometer-last-values
    total_transactions_last_values = total_transactions_element.find_all('div', class_='odometer-last-value')
    total_transactions_integers = ''.join(value.get_text(strip=True) for value in total_transactions_last_values if value.get_text(strip=True).isdigit())
    try:
        total_transactions_integers = format_number(int(total_transactions_integers))
    except:
        error_message = ttk.Label(master=window,
                                text="Die gefundenen Informationen sind nicht vorhanden."
                                "Überprüfe den Kanal und die angegebenen Analyse-Daten",
                                font= "Calibri 28 bold",
                                foreground='red')
        error_message.grid(row=5, column=0, pady=15)

    # Extract total sum region
    total_sum_region = soup.find('span', id='total_sum')

    # Extract total sum odometer-last-values
    total_sum_last_values = total_sum_region.find_all('div', class_='odometer-last-value')
    total_sum_integers = ''.join(value.get_text(strip=True) for value in total_sum_last_values if value.get_text(strip=True).isdigit())

    markings = total_sum_region.find_all('span', class_="odometer-formatting-mark")
    
    try:
        total_sum_integers = format_money(int(total_sum_integers), markings)
    except:
        pass
        
    # Print the results
    """     print(f"The total transactions are: {total_transactions_integers}")
    print(f"The total sum is: {total_sum_integers} EUR") """
    return total_transactions_integers, total_sum_integers
def extract_geographic(html_text):
    # Create a BeautifulSoup object with the HTML text
    soup = BeautifulSoup(html_text, 'html.parser')

    # Find the table containing the country data
    country_table = soup.find('table', class_='table-striped')

    # Initialize an empty dictionary to store the country data
    country_data = {}

    count = 0

    # Extract data from each row in the table
    for row in country_table.find_all('tr') [1:]:  # Skip the header row
        columns = row.find_all('td')
    
        # Extract country name and unique views
        country_name = columns[0].get_text(strip=True)
        unique_views = int(columns[1].get_text(strip=True))
        unique_views = format_number(unique_views)
    
        # Store data in the dictionary
        country_data[country_name] = unique_views

        count += 1
        if count == 5:
            break

    # Print the outcome
    """     print("The top 5 countries people viewed from by unique views are:")
    for country, views in country_data.items():
        print(f"{country}: {views} unique views") """
    return country_data
def extract_toplist_livestreams(html_text):
    soup = BeautifulSoup(html_text, 'html.parser')

    # Find the table containing the country data
    toplist_livestreams_table = soup.find('table', class_='table-striped')

    # Initialize an empty dictionary to store the country data
    toplist_livestreams_data = {}

    # Counter variable to only list the first five entries
    count = 0

    # Extract data from each row in the table
    for row in toplist_livestreams_table.find_all('tr')[1:]:  # Skip the header row
        columns = row.find_all('td')
        
        # Extract toplist_livestream and unique views
        toplist_livestreams = columns[0].get_text(strip=True)
        unique_views = int(columns[1].get_text(strip=True))
        unique_views = format_number(unique_views)
        
        # Store data in the dictionary
        toplist_livestreams_data[toplist_livestreams] = unique_views

        count += 1 
        
        # Check if we have processed five entries, and if so, break out of the loop
        if count == 5:
            break

    # Print the outcome
    """ print("Top 5 livestreams by unique views:\n")
    for livestream, views in toplist_livestreams_data.items():
        print(f"{livestream}: {views} unique views") """
    return toplist_livestreams_data
def extract_toplist_videos(html_text):
    # Create a BeautifulSoup object with the HTML text
    soup = BeautifulSoup(html_text, 'html.parser')

    # Find the table containing the country data
    toplist_videos_table = soup.find('table', class_='table-striped')

    # Initialize an empty dictionary to store the country data
    toplist_videos_data = {}

    # Counter variable to only list the first five entries
    count = 0

    # Extract data from each row in the table
    for row in toplist_videos_table.find_all('tr')[1:]:  # Skip the header row
        columns = row.find_all('td')
        
        # Extract toplist_video and unique views
        toplist_videos = columns[0].get_text(strip=True)
        unique_views = int(columns[1].get_text(strip=True))
        unique_views = format_number(unique_views)
        
        # Store data in the dictionary
        toplist_videos_data[toplist_videos] = unique_views

        count += 1 
        
        # Check if we have processed five entries, and if so, break out of the loop
        if count == 5:
            break

    # Print the outcome
    """ print("Top 5 videos by unique views:\n")
    for video, views in toplist_videos_data.items():
        print(f"{video}: {views} unique views") """
    return toplist_videos_data
def extract_toplist_categories(html_text):
    # Create a BeautifulSoup object with the HTML text
    soup = BeautifulSoup(html_text, 'html.parser')

    # Find the table containing the country data
    toplist_categories_table = soup.find('table', class_='table-striped')

    # Initialize an empty dictionary to store the country data
    toplist_categories_data = {}

    # Counter variable to only list the first five entries
    count = 0

    # Extract data from each row in the table
    for row in toplist_categories_table.find_all('tr')[1:]:  # Skip the header row
        columns = row.find_all('td')
        
        # Extract toplist_categorie and unique views
        toplist_categories = columns[0].get_text(strip=True)
        unique_views = int(columns[1].get_text(strip=True))
        unique_views = format_number(unique_views)
        
        # Store data in the dictionary
        toplist_categories_data[toplist_categories] = unique_views

        count += 1 
        
        # Check if we have processed five entries, and if so, break out of the loop
        if count == 5:
            break

    # Print the outcome
    """ print("Top 5 categories by unique views:\n")
    for categorie, views in toplist_categories_data.items():
        print(f"{categorie}: {views} unique views") """
    return toplist_categories_data
def extract_toplist_downloads(html_text):
    # Create a BeautifulSoup object with the HTML text
    soup = BeautifulSoup(html_text, 'html.parser')

    # Find the table containing the country data
    toplist_downloads_table = soup.find('table', class_='table-striped')

    # Initialize an empty dictionary to store the country data
    toplist_downloads_data = {}

    # Counter variable to only list the first five entries
    count = 0

    # Extract data from each row in the table
    for row in toplist_downloads_table.find_all('tr')[1:]:  # Skip the header row
        columns = row.find_all('td')
        
        # Extract toplist_download and unique views
        toplist_downloads = columns[0].get_text(strip=True)
        unique_views = int(columns[1].get_text(strip=True))
        unique_views = format_number(unique_views)
        
        # Store data in the dictionary
        toplist_downloads_data[toplist_downloads] = unique_views

        count += 1 
        
        # Check if we have processed five entries, and if so, break out of the loop
        if count == 5:
            break  
    # Print the outcome
    """ print(f"Top {count} downloads by downloads:\n")
    for download, views in toplist_downloads_data.items():
        print(f"{download}: {views} downloads") """
    return toplist_downloads_data
def extract_all_access_data(html_text):
    soup = BeautifulSoup(html_text, 'html.parser')

    all_access_passes = []
    all_access_data = {}

    keys = soup.find_all('th')
    values = soup.find_all('td')

    count = 0
    index = []
    for value in values:
        if str(value.get_text(strip=True)).isdigit():
            if int(value.get_text(strip=True)) != 0:
                index.append(count)
        count += 1
    
    for place in index:
        all_access_data[keys[1].get_text(strip=True)] = values[place-5].get_text(strip=True)
        all_access_data[keys[2].get_text(strip=True)] = values[place-4].get_text(strip=True)
        all_access_data[keys[6].get_text(strip=True)] = values[place].get_text(strip=True)
        all_access_passes.append(all_access_data.copy())

    print(all_access_passes)
    return all_access_passes
def extract_billing_data(html_text):
    soup = BeautifulSoup(html_text, 'html.parser')

    billing_info = {'Pay Per View reimbursement up until': '',
                    'reference_number': '',
                    'amount': '',
                    'conditions_of_payment': ''}
    count = 1
    # Find the invoice details
    invoice_details = soup.find('div', id='currency_price_EUR')

    if invoice_details:
        # Extract invoice information
        invoice_info = invoice_details.find_all('strong')
        if invoice_info:
            # Iterate over invoice_info starting from the second item
            for item_key, item_value in billing_info.items():
                if count >= len(invoice_info):
                    break  # Break if no more items in invoice_info
                if count == 0:
                    count += 1  # Skip the first item
                    continue
                item = invoice_info[count]  # Get the corresponding item from invoice_info
                billing_info[item_key] = item.text.strip()
                count += 1
    print(billing_info)  
    return billing_info  
def extract_multi_currency_data(html_text):
    soup = BeautifulSoup(html_text, 'html.parser')
    multi_currency_amounts = []

    currency_details = soup.find(id='currency_price_tabs')
    count = 0
    for currency in currency_details:
        if (count % 2):
            multi_currency_amounts.append(currency.get_text(strip=True))
            count += 1
        else:
            count += 1

    print(f'!!!!! {multi_currency_amounts}')
    return multi_currency_amounts
def extract_selltements_data(html_text):
    soup = BeautifulSoup(html_text, 'html.parser')
    settlements = {}

    keys = soup.find_all('th')
    values = soup.find_all('td')

    for key in keys:
        print(f'KEY: {key.get_text(strip=True)}')
    for value in values:
        print(f'VALUE: {value.get_text(strip=True)}')

    for key, value in zip(keys, values):
        key_text = key.get_text(strip=True)
        value_text = value.get_text(strip=True)
        settlements[key_text] = value_text

    return settlements
def transactions_read(user_info_dictionary):
    # Specify the path to your text document on the desktop
    desktop_path = f'~/Desktop/{user_info_dictionary['event_name']}'

    # Expand the ~ character to the home directory
    desktop_path = os.path.expanduser(desktop_path)

    # List all files in the directory
    files_in_directory = os.listdir(desktop_path)

    # Filter out the CSV files
    csv_files = [file for file in files_in_directory if file.endswith('.csv')]

    if len(csv_files) != 1:
        print("Error: There should be exactly one CSV file in the directory.")
        exit()

    csv_filename = os.path.join(desktop_path, csv_files[0])

    data_category = 'All-access'
    html_text = access_file(data_category, user_info_dictionary)
    all_access_data = extract_all_access_data(html_text)

    # Open the CSV file
    with open(f'{csv_filename}', 'r', newline='', encoding='utf-8') as file:
        # Create a CSV reader object
        csv_reader = csv.DictReader(file)
        
        # Skip the header row
        next(csv_reader)
        
        # Iterate through each row
        transaction_count = 0
        transaction_currency_count = [0, 0, 0, 0]
        singel_stream_currency_count = [0, 0, 0, 0]
        brut_rev = [0, 0, 0, 0]
        single_stream_prices = {'SEK': '-', 'EUR': '-', 'USD': '-', 'NOK': '-'}
        object_prices = {}
        object_purchase_counts = {} 
        for row in csv_reader:
            # Increment transaction count
            transaction_count += 1
            
            # Extract data from the row
            object_name = row['Object']
            amount = row['Amount']

            if 'SEK' in amount:
                transaction_currency_count[0] += 1
            elif 'EUR' in amount:
                transaction_currency_count[1] += 1
            elif 'USD' in amount:
                transaction_currency_count[2] += 1
            elif 'NOK'in amount:
                transaction_currency_count[3] += 1

            for all_access_passes in all_access_data:
                if str(all_access_passes['Title'][:20]) not in object_name[:20]:
                    if 'SEK' in amount:
                        singel_stream_currency_count[0] += 1
                    elif 'EUR' in amount:
                        singel_stream_currency_count[1] += 1
                    elif 'USD' in amount:
                        singel_stream_currency_count[2] += 1
                    elif 'NOK'in amount:
                        singel_stream_currency_count[3] += 1
            
            # split the amount String for object_purchase_counts
            price, currency = amount.split()

            if 'SEK' in currency:
                brut_rev[0] += float(price)
            elif 'EUR' in currency:
                brut_rev[1] += float(price)
            elif 'USD' in currency:
                brut_rev[2] += float(price)
            elif 'NOK'in currency:
                brut_rev[3] += float(price)

            # Update object_prices dictionary
            if object_name in object_prices:
                if amount not in object_prices[object_name]:
                    object_prices[object_name].append(amount)
            else:
                object_prices[object_name] = [amount]
            
            # Update object_purchase_counts dictionary
            
            # Update object_purchase_counts dictionary
            if object_name in object_purchase_counts:
                if currency in object_purchase_counts[object_name]:
                    object_purchase_counts[object_name][currency] += 1
                else:
                    object_purchase_counts[object_name][currency] = 1
            else:
                object_purchase_counts[object_name] = {currency: 1}

            bought_all_access_passes = {}
        for all_access_passes in all_access_data:
            for object in object_purchase_counts:
                if str(all_access_passes['Title'][:26]) in object[:26]:
                    bought_all_access_passes[all_access_passes['Title']] = []
                    bought_all_access_passes[all_access_passes['Title']].append(
                        object_prices[object]
                    )
                    bought_all_access_passes[all_access_passes['Title']].append(
                        object_purchase_counts[object]
                    )
        count = 0
        for object in object_prices:
            if str(all_access_passes['Title'][:26]) not in object[:26]:
                if 'SEK' in object_prices[object][count]:
                    single_stream_prices['SEK'] = (object_prices[object][count]
                                                   .replace('SEK', ''))
                elif 'EUR' in object_prices[object][count]:
                    single_stream_prices['EUR'] = (object_prices[object][count]
                                                   .replace('EUR', ''))
                elif 'USD' in object_prices[object][count]:
                    single_stream_prices['USD'] = (object_prices[object][count]
                                                   .replace('USD', ''))
                elif 'NOK' in object_prices[object][count]:
                    single_stream_prices['NOK'] = (object_prices[object][count]
                                                   .replace('NOK', ''))

        print(bought_all_access_passes)
        print(transaction_currency_count)
        print(f'!!! {singel_stream_currency_count}')
        print(single_stream_prices)

    return (transaction_count, bought_all_access_passes, 
            transaction_currency_count, single_stream_prices, brut_rev, 
            singel_stream_currency_count)

# Powerpoint Table functions
def create_table(slide, rows_count, columns_count):
    table_placeholder = slide
    graphic_frame = table_placeholder.insert_table(rows=rows_count, cols=columns_count)
    table = graphic_frame.table 

    tbl =  graphic_frame._element.graphic.graphicData.tbl
    style_id = '{F2DE63D5-997A-4646-A377-4702673A728D}'
    tbl[0][-1].text = style_id
    return table
def access_and_fill_table(table, info_list):
    row = 1
    column = 0
    for key, value in info_list.items():
        if row >= len(table.rows):
            break  # Exit the loop if we reach the end of the table

        cell_key = table.cell(row, column)
        cell_value = table.cell(row, column + 1)

        cell_key.text = key
        cell_value.text = str(value)

        row += 1  # Move to the next row

        # Reset column for the next row (assuming you want to start with column 0)
        column = 0
def make_text_bold_in_table(table, text_to_bold):
    for row in table.rows:
        for cell in row.cells:
            if cell.text in text_to_bold:
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text == cell.text:
                                run.font.bold = True
                                run.font.color.rgb = RGBColor(225, 39, 42)
            else:
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(16)
def create_revenue_table_1(slide, rows_count, columns_count, 
                         all_access_data):
    table_placeholder = slide
    graphic_frame = table_placeholder.insert_table(
        rows=rows_count+((len(all_access_data)-1)*2), cols=columns_count)
    table = graphic_frame.table 
    cell1 = table.cell(1,0)
    cell2 = table.cell(2+(len(all_access_data)-1),0)
    cell1.merge(cell2)

    cell3 = table.cell(3+(len(all_access_data)-1),0)
    cell4 = table.cell(5+((len(all_access_data)-1)*2),0)
    cell3.merge(cell4)
     

    tbl =  graphic_frame._element.graphic.graphicData.tbl
    style_id = '{5940675A-B579-460E-94D1-54222C63F5DA}'
    tbl[0][-1].text = style_id
    return table
def fill_reavenue_table_1(table, user_info_dictionary, ppv_info_list, 
                        billing_info, all_access_data, 
                        single_stream_purchase_count, multi_currency_info,):
    cell_data = {
        (0, 0): 'Categorie',
        (0, 1): 'Description',
        (0, 2): 'TOTAL',
        (1, 0): 'Prices',
        (1, 1): 'Price Singlestream',
        (1, 2): f"{str(user_info_dictionary['price_single_stream'])} EUR",
        (2, 1): 'Price All Access Pass',
        (2, 2): str(all_access_data['all_access_price']),
        (3, 0): 'Transactions',
        (3, 1): 'Total transactions',
        (3, 2): str(ppv_info_list[0]),
        (4, 1): 'Total transactions single streams',
        (4, 2): str(single_stream_purchase_count),
        (5, 1): 'Total transactions All Access Pass',
        (5, 2): f'{str(all_access_data[0]['Active subscribers'])}',
        (6, 0): 'Total reavenue*',
        (6, 2): f'{str(ppv_info_list[1])} EUR',
        (7, 0): 'Settlements',
        (7, 2): str(user_info_dictionary['revenue_split']),
        (8, 0): 'Net reavenue*',
        (8, 2): f'{str(billing_info['amount'])}',
    }
    for (row, col), text in cell_data.items():
        cell = table.cell(row, col)
        cell.text = text
def create_revenue_table_2(slide, rows_count, columns_count, all_access_data,
                            settlements_data):
    table_placeholder = slide
    graphic_frame = table_placeholder.insert_table(
        rows=rows_count, cols=(columns_count+((len(all_access_data)-1)*2)
        +(len(settlements_data)-1)))
    table = graphic_frame.table 
    
    cell1 = table.cell(0,1)
    cell2 = table.cell(0,2+(len(all_access_data)-1))
    cell1.merge(cell2)

    cell3 = table.cell(0,3+(len(all_access_data)-1))
    cell4 = table.cell(0,5+((len(all_access_data)-1)*2))
    cell3.merge(cell4)

    if len(settlements_data) > 1:
        cell5 = table.cell(0, 7+((len(all_access_data)-1)*2))
        cell6 = table.cell(0, 7+((len(all_access_data)-1)*2)
                           +(len(settlements_data)-1))
        cell5.merge(cell6)

    cell7 = table.cell(2, 0)
    cell10 = table.cell(5, 0)
    #cell8.merge(cell7)
    cell10.merge(cell7)

    tbl =  graphic_frame._element.graphic.graphicData.tbl
    style_id = '{5940675A-B579-460E-94D1-54222C63F5DA}'
    tbl[0][-1].text = style_id
    return table
def fill_reavenue_table_2(table, user_info_dictionary, ppv_info_list, 
                        billing_info, all_access_data, 
                        single_stream_purchase_count, multi_currency_info,
                        settlements_data, transaction_data):
    (transaction_count, bought_all_access_passes, transaction_currency_count, 
    single_stream_prices, brut_rev, singel_stream_currency_count) = transaction_data
    cell_data = {
        (0, 0): 'Categorie',
        (1, 0): 'Description',
        (2, 0): 'Total',
        (6, 0): '∑',
        (0, 1): 'Prices',
        (0, 3+(len(all_access_data)-1)): 'Transactions',
        (0, 6+(len(all_access_data)-1)*2): 'Total gross*',
        (0, 7+(len(all_access_data)-1)*2): 'Settlements',
        (0, 8+((len(all_access_data)-1)*2)+(len(settlements_data)-1)):
        'Net reavenue*',
        (1, 1): 'Preis singelstream',
        (2, 1): single_stream_prices['SEK'] + ' SEK',
        (3, 1): single_stream_prices['EUR'] + ' EUR',
        (4, 1): single_stream_prices['USD'] + ' USD',
        (5, 1): single_stream_prices['NOK'] + ' NOK',
        (6, 5+(len(all_access_data)-1)*2): str(transaction_count),
        (1, 3+(len(all_access_data)-1)): 'Transactions singlestream',
        (2, 3+(len(all_access_data)-1)): str(singel_stream_currency_count[0]),
        (3, 3+(len(all_access_data)-1)): str(singel_stream_currency_count[1]),
        (4, 3+(len(all_access_data)-1)): str(singel_stream_currency_count[2]),
        (5, 3+(len(all_access_data)-1)): str(singel_stream_currency_count[3]),
        (1, 5+(len(all_access_data)-1)*2): 'Total transactions',
        (2, 5+(len(all_access_data)-1)*2): str(transaction_currency_count[0]),
        (3, 5+(len(all_access_data)-1)*2): str(transaction_currency_count[1]),
        (4, 5+(len(all_access_data)-1)*2): str(transaction_currency_count[2]),
        (5, 5+(len(all_access_data)-1)*2): str(transaction_currency_count[3]),
        (2, 6+(len(all_access_data)-1)*2): str(round(brut_rev[0], 2)) + ' SEK',
        (3, 6+(len(all_access_data)-1)*2): str(round(brut_rev[1], 2)) + ' EUR',
        (4, 6+(len(all_access_data)-1)*2): str(round(brut_rev[2], 2)) + ' USD',
        (5, 6+(len(all_access_data)-1)*2): str(round(brut_rev[3], 2)) + ' NOK',
        (2, 8+(len(all_access_data)-1)*2): str(multi_currency_info[0]),
        (3, 8+(len(all_access_data)-1)*2): str(multi_currency_info[0]),
        (4, 8+(len(all_access_data)-1)*2): str(multi_currency_info[0]),
        (5, 8+(len(all_access_data)-1)*2): str(multi_currency_info[0])
    }

    for index, data in enumerate(bought_all_access_passes, start=2):
        cell_data[(1, index)] = data
        for prices in bought_all_access_passes[data][0]:
            if 'SEK' in prices:
                cell_data[(2, index)] = str(prices)
            elif 'EUR' in prices:
                cell_data[(3, index)] = str(prices)
            elif 'USD' in prices:
                cell_data[(4, index)] = str(prices)
            else:
                cell_data[(5, index)] = str(prices)
        cell_data[(1, 2+index+len(all_access_data)-1)] = data
        cell_data[(2, 2+index+len(all_access_data)-1)] = str(
        bought_all_access_passes[data][1].get('SEK', '0'))
        cell_data[(3, 2+index+len(all_access_data)-1)] = str(
            bought_all_access_passes[data][1].get('EUR', '0'))
        cell_data[(4, 2+index+len(all_access_data)-1)] = str(
            bought_all_access_passes[data][1].get('USD', '0'))
        cell_data[(5, 2+index+len(all_access_data)-1)] = str(
            bought_all_access_passes[data][1].get('NOK', '0'))
        cell_data[(6, 2+index+len(all_access_data)-1)] = all_access_data[
            index-2]['Active subscribers']

    print(f'{transaction_count} \n')
    print(f'{bought_all_access_passes} \n')
    print(f'{transaction_currency_count} \n')
    print(f'{single_stream_prices} \n')

    for (row, col), text in cell_data.items():
        cell = table.cell(row, col)
        cell.text = text
def format_table(table, text_to_bold):
    for row in table.rows:
        for cell in row.cells:
            if cell.text in text_to_bold:
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text == cell.text:
                                run.font.bold = True
                                run.font.color.rgb = RGBColor(225, 39, 42)
                                run.font.size = Pt(12)
            else:
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(11)                                
def create_analytics_table_1(slide, rows_count, columns_count):
    table_placeholder = slide
    graphic_frame = table_placeholder.insert_table(rows=rows_count, cols=columns_count)
    table = graphic_frame.table 

    tbl =  graphic_frame._element.graphic.graphicData.tbl
    style_id = '{5940675A-B579-460E-94D1-54222C63F5DA}'
    tbl[0][-1].text = style_id
    return table
def format_analytics_table(table, text):
    for row in table.rows:
        for cell in row.cells:
            if cell.text in text:
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text == cell.text:
                                run.font.bold = True
                                run.font.color.rgb = RGBColor(225, 39, 42)
                                run.font.size = Pt(10)
            else:
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(9)    

# Main function
def main():
    user_info_dictionary = {'event_name': '',
                        'page_url': '',
                          'e_mail': '',
                            'password': '',
                              'start_date': '',
                                'end_date':'',
                                'revenue_split': '',
                                'price_single_stream': '',
                                'all_access_pass_exist': '',
                                'multi_currrency': '',
                                  'categories_exist': '',
                                    'downloads_exist': '',
                                    'powered_by': ''}
    user_info_list = retrieve_info()
    print(user_info_list)

    # Unpack user_info_list into corresponding dictionary values
    (user_info_dictionary['event_name'],
    user_info_dictionary['page_url'],
    user_info_dictionary['e_mail'],
        user_info_dictionary['password'],
        user_info_dictionary['start_date'],
            user_info_dictionary['end_date'],
            user_info_dictionary['revenue_split'],
            user_info_dictionary['price_single_stream'],
            user_info_dictionary['all_access_pass_exist'],
            user_info_dictionary['multi_currrency'],
            user_info_dictionary['categories_exist'],
                user_info_dictionary['downloads_exist'],
                user_info_dictionary['powered_by']) = user_info_list

    print(user_info_dictionary)

    folder_creation(user_info_dictionary)

    # Initialize driver
    folder_path = os.path.expanduser(f"~/Desktop/{user_info_dictionary['event_name']}")
    driver = initialize_driver(folder_path)

    # Update the link to the new login page
    login_url = "https://www.solidsport.com/login"

    # Perform the login
    driver.get(login_url)

    # Handle cookie popup (if present)
    handle_cookie_popup(driver)

    # Perform the login
    login(driver, user_info_dictionary)

    # Navigate to the Winterturnier page
    navigate_to_page(driver, user_info_dictionary)

    if user_info_dictionary['powered_by'] == True:
        powered_by_navigation(driver)
    else:
        # Click on the "Admin" button
        admin_button = WebDriverWait(driver, 3).until(
            EC.element_to_be_clickable((By.XPATH, "//div[text()='Admin']"))
        )
        admin_button.click()

        # Click on the "Analytics" button
        analytics_button = WebDriverWait(driver, 1).until(
            EC.element_to_be_clickable((By.XPATH, "//span[text()='Analytics']"))
        )
        analytics_button.click()

    # Enter dates and submit
    enter_dates_and_submit(driver, user_info_dictionary)

    # Funciton capture_information_overwiew
    capture_information(driver, "overview", user_info_dictionary)

     # Click on the "Geographic" button
    geographic_button = WebDriverWait(driver, 30).until(
    EC.element_to_be_clickable((By.LINK_TEXT, 'Geographic'))
    )   
    geographic_button.click()

    #time.sleep(2)
    capture_information(driver, "geographic", user_info_dictionary)

    # Click on the "Toplist" button
    toplist_button = WebDriverWait(driver, 30).until(
    EC.element_to_be_clickable((By.LINK_TEXT, 'Toplist'))
    )   
    toplist_button.click()

    #time.sleep(2)
    capture_information(driver, "toplist_videos", user_info_dictionary)

    # Click on the "Livestreams" button
    toplist_button = WebDriverWait(driver, 30).until(
    EC.element_to_be_clickable((By.LINK_TEXT, 'Livestreams'))
    )   
    toplist_button.click()

    time.sleep(2)
    capture_information(driver, "toplist_livestreams", user_info_dictionary)

    # Click on the "Categories" button
    if user_info_dictionary['categories_exist'] == True:
        toplist_button = WebDriverWait(driver, 30).until(
        EC.element_to_be_clickable((By.LINK_TEXT, 'Categories'))
        )   
        toplist_button.click()
    capture_information(driver, "toplist_categories", user_info_dictionary)

    # Click on the "Downloads" button
    if user_info_dictionary['downloads_exist'] == True:
        toplist_button = WebDriverWait(driver, 30).until(
        EC.element_to_be_clickable((By.LINK_TEXT, 'Downloads'))
        )   
        toplist_button.click()
    capture_information(driver, "toplist_downloads", user_info_dictionary)

    # Click on the "ppv" button
    ppv_button = WebDriverWait(driver, 30).until(
    EC.element_to_be_clickable((By.LINK_TEXT, 'PPV'))
    )
    ppv_button.click()

    time.sleep(2)
    capture_information(driver, "ppv", user_info_dictionary)

     # Click on the "Shop" button
    shop_button = WebDriverWait(driver, 1).until(
        EC.element_to_be_clickable((By.LINK_TEXT, 'Shop'))
    )
    shop_button.click()
    download_transaction_data(driver, user_info_dictionary['start_date'],
                              user_info_dictionary['end_date'], 60)
    time.sleep(10)
    # Click on the "All-Access" button
    all_access_button = WebDriverWait(driver, 1).until(
        EC.element_to_be_clickable((By.LINK_TEXT, 'All-access'))
    )
    all_access_button.click()
    capture_information(driver, 'All-access', user_info_dictionary)

    # Click on the "Billing" button
    billing_button = WebDriverWait(driver, 1).until(
        EC.element_to_be_clickable((By.LINK_TEXT, 'Billing'))
    )
    billing_button.click()
    capture_information(driver, 'billing', user_info_dictionary)

    # Click on the 'Settlements' button
    settlement_button = WebDriverWait(driver, 1).until(
        EC.element_to_be_clickable((By.LINK_TEXT, 'Settlements'))
    )
    settlement_button.click()
    capture_information(driver, 'settlements', user_info_dictionary)

    """ # Change the language to german
    german_flag_button = WebDriverWait(driver, 1).until(
        EC.element_to_be_clickable((By.CLASS_NAME, 'flag-de'))
    )
    german_flag_button.click() """
    
    # Close the browser
    driver.quit()

    data_category = "overview"
    html_text = access_file(data_category, user_info_dictionary)
    overview_info_list = extract_views(html_text)

    data_category = "ppv"
    html_text = access_file(data_category, user_info_dictionary)
    ppv_info_list = extract_ppv_info(html_text)

    data_category = "geographic"
    html_text = access_file(data_category, user_info_dictionary)
    geographic_info_list = extract_geographic(html_text)

    data_category = "toplist_livestreams"
    html_text = access_file(data_category, user_info_dictionary)
    toplist_livestreams_info_list = extract_toplist_livestreams(html_text)

    data_category = "toplist_videos"
    html_text = access_file(data_category, user_info_dictionary)
    toplist_videos_info_list = extract_toplist_videos(html_text)

    if user_info_dictionary['categories_exist'] == True:
        data_category = "toplist_categories"
        html_text = access_file(data_category, user_info_dictionary)
        toplist_categories_info_list = extract_toplist_categories(html_text)
        print(toplist_categories_info_list)

    if user_info_dictionary['downloads_exist'] == True:
        data_category = "toplist_downloads"
        html_text = access_file(data_category, user_info_dictionary)
        toplist_downloads_info_list = extract_toplist_downloads(html_text)
        print(toplist_downloads_info_list)

    data_category = 'All-access'
    html_text = access_file(data_category, user_info_dictionary)
    all_access_data = extract_all_access_data(html_text)

    data_category = 'billing'
    html_text = access_file(data_category, user_info_dictionary)
    billing_info = extract_billing_data(html_text)
    multi_currency_info = extract_multi_currency_data(html_text)

    data_category = 'billing'
    html_text = access_file(data_category, user_info_dictionary)
    settlements_data = extract_selltements_data(html_text)
    settlements_data = ['0']       

    desktop_path = os.path.join(os.path.expanduser('~'), 'Desktop')
    file_name = 'reportvorlage_EN.pptx'
    file_path = os.path.join(desktop_path, file_name)
    
    total_all_access_passes = calculate_all_bought_all_access_passes(all_access_data)
    print(f'Total_all_access_passes: {total_all_access_passes}')

    single_stream_purchase_count = calculate_single_stream_purchase_count(
        ppv_info_list[0], total_all_access_passes)
    
    calculations = []
    calculations = calculate_analytic_values(overview_info_list[1],
                                             overview_info_list[0],
                                             ppv_info_list[0], 
                                             total_all_access_passes,
                                             (single_stream_purchase_count),
                                             ppv_info_list[1],
                                             billing_info['amount'], 
                                             user_info_dictionary['revenue_split'])

    transaction_data = transactions_read(user_info_dictionary)

    # create Powerpoint
    prs = Presentation(file_path)

    #register slides
    title_slide_register = prs.slide_layouts[0]
    overview_slide_register = prs.slide_layouts[1]
    toplist_streams_slide_register = prs.slide_layouts[2]
    toplist_downloads_slide_register = prs.slide_layouts[3]
    seciton_earnings_slide_register = prs.slide_layouts[4]
    overview_earnings_slide_register = prs.slide_layouts[5]
    multi_currency_slide_register = prs.slide_layouts[6]
    billing_slide_register = prs.slide_layouts[7]
    further_info_slide_register = prs.slide_layouts[8]

    #Add slides to presentation
    title_slide = prs.slides.add_slide(title_slide_register)
    overview_slide = prs.slides.add_slide(overview_slide_register)
    toplist_streams_slide = prs.slides.add_slide(toplist_streams_slide_register)
    if user_info_dictionary['downloads_exist'] or user_info_dictionary['categories_exist'] == True:
        toplist_downloads_slide = prs.slides.add_slide(toplist_downloads_slide_register)
    seciton_earnings_slide = prs.slides.add_slide(seciton_earnings_slide_register)
    if (user_info_dictionary['multi_currrency'] == False and len(all_access_data) == 1 and 
    len(settlements_data) == 1):
        overview_earnings_slide = prs.slides.add_slide(overview_earnings_slide_register)
    else:
        multi_currency_slide = prs.slides.add_slide(multi_currency_slide_register)
    billing_slide = prs.slides.add_slide(billing_slide_register)
    further_info_slide = prs.slides.add_slide(further_info_slide_register)

    # title slide
    title_slide.placeholders[0].text = user_info_dictionary['event_name']
    title_slide.placeholders[1].text = f"Summary, analysis, overview of earnings, and invoicing: {user_info_dictionary['event_name']}"
    title_slide.placeholders[14].text = 'Quelle: Solidsport'

    # overview slide
    # Add each bullet point as a separate paragraph
    bullet_points = [
        f"This event had a total of: {overview_info_list[0]} views",
        f"of those: {overview_info_list[1]} were unique views",
        f"This event had a total of: {ppv_info_list[0]} transaktions",
        f"The gross reavenue: {ppv_info_list[1]} ",
        f"You net share: {billing_info['amount']}"
    ]

    # Add each bullet point
    for point in bullet_points:
        p = overview_slide.placeholders[2].text_frame.add_paragraph()
        p.text = point
        p.level = 0  # Set the indentation level for bullet points
        p.space_before = Pt(0)  # Set space before paragraph to zero
    
    overview_slide.placeholders[3].text = f'Top {len(geographic_info_list)} countries according to unique views'
    
    geographic_table = create_table(overview_slide.placeholders[15], len(geographic_info_list), 2)
    header_cell_1 = geographic_table.cell(0,0)
    header_cell_1.text = "Countries"
    header_cell_2 = geographic_table.cell(0,1)
    header_cell_2.text = "unique views"
    access_and_fill_table(geographic_table, geographic_info_list)
    make_text_bold_in_table(geographic_table, ["Countries", 'unique views'])

    analytics_table = create_analytics_table_1(overview_slide.placeholders[16], 3, 2)
    header_cell_analyse_1 = analytics_table.cell(0,0)
    header_cell_analyse_1.text = "unique views / views"
    header_cell_analyse_2 = analytics_table.cell(0,1)
    header_cell_analyse_2.text = "unique views / transactions"

    analyse_cell_10 = analytics_table.cell(1,0)
    analyse_cell_10.text = f"{overview_info_list[1]} / {overview_info_list[0]} ≈ {calculations[0]*100}%"
    analyse_cell_20 = analytics_table.cell(2,0)

    if calculations[0] >= 0.74:
        analyse_cell_20.text = f'A ratio of: {calculations[0]*100}% is very good!'
    else:
        analyse_cell_20.text = f'A ratio of: {calculations[0]*100}% is avaage!'

    analyse_cell_11 = analytics_table.cell(1,1)
    analyse_cell_11.text = f"{overview_info_list[1]} / {ppv_info_list[0]} ≈ {calculations[1]}"
    analyse_cell_21 = analytics_table.cell(2,1)
    analyse_cell_21.text = f"Per purchase a person viewed ⌀ {calculations[1]} games"

    format_analytics_table(analytics_table, ["Einmalige Ansichten / Ansichten", 'Einmalige Ansichten / Transaktionen'])

    overview_slide.placeholders[13].text = f'unique views: from({user_info_dictionary['start_date']}) | til({user_info_dictionary['end_date']})'
    overview_slide.placeholders[14].text = 'Source: Solidsport'

    # toplist streams
    toplist_streams_slide.placeholders[1].text = f'Top {len(toplist_livestreams_info_list)} Livestreams pre unique views'

    toplist_livestreams_table = create_table(toplist_streams_slide.placeholders[16], len(toplist_livestreams_info_list)+1, 2)
    header_cell_1 = toplist_livestreams_table.cell(0,0)
    header_cell_1.text = "Livestream"
    header_cell_2 = toplist_livestreams_table.cell(0,1)
    header_cell_2.text = "unique views"
    access_and_fill_table(toplist_livestreams_table, toplist_livestreams_info_list)
    make_text_bold_in_table(toplist_livestreams_table, ['Livestream', 'unique views'])

    toplist_streams_slide.placeholders[3].text = f'Top {len(toplist_videos_info_list)} Videos per unique views'
    
    toplist_videos_table = create_table(toplist_streams_slide.placeholders[15], len(toplist_videos_info_list)+1, 2)
    header_cell_1 = toplist_videos_table.cell(0,0)
    header_cell_1.text = "Videos"
    header_cell_2 = toplist_videos_table.cell(0,1)
    header_cell_2.text = "unique views"
    access_and_fill_table(toplist_videos_table, toplist_videos_info_list)
    make_text_bold_in_table(toplist_videos_table, ['Videos', 'unique views'])

    toplist_streams_slide.placeholders[13].text = f'unique views: from({user_info_dictionary['start_date']}) | till({user_info_dictionary['end_date']})'
    toplist_streams_slide.placeholders[14].text = 'source: Solidsport'

    # toplist downloads slide
    if user_info_dictionary['downloads_exist'] and user_info_dictionary['categories_exist'] == True:
        #toplist_downloads_slide.placeholders[0].text = 'Beliebteste Kategorien und Top Downloads nach einmaligen Ansichten'
        toplist_downloads_slide.placeholders[1].text = f'Top {len(toplist_categories_info_list)} Categories per unique views'
        toplist_categories_table = create_table(toplist_downloads_slide.placeholders[16], len(toplist_categories_info_list)+1, 2)
        header_cell_1 = toplist_categories_table.cell(0,0)
        header_cell_1.text = "Categorie"
        header_cell_2 = toplist_categories_table.cell(0,1)
        header_cell_2.text = "unique views"
        access_and_fill_table(toplist_categories_table, toplist_categories_info_list)
        make_text_bold_in_table(toplist_categories_table, ['Categorie', 'unique views'])
        toplist_downloads_slide.placeholders[3].text = f'The {len(toplist_downloads_info_list)} most downloaded videos'
        toplist_downloads_table = create_table(toplist_downloads_slide.placeholders[15], len(toplist_downloads_info_list)+1, 2)
        header_cell_1 = toplist_downloads_table.cell(0,0)
        header_cell_1.text = "Video"
        header_cell_2 = toplist_downloads_table.cell(0,1)
        header_cell_2.text = "Downloads"
        access_and_fill_table(toplist_downloads_table, toplist_downloads_info_list)
        make_text_bold_in_table(toplist_downloads_table, ['Video', 'Downloads'])
        toplist_downloads_slide.placeholders[13].text = f'unique views: from({user_info_dictionary['start_date']}) | til({user_info_dictionary['end_date']})'
        toplist_downloads_slide.placeholders[14].text = 'Source: Solidsport'
    elif user_info_dictionary['downloads_exist'] == True:
        #toplist_downloads_slide.placeholders[0].text = 'Top Downloads nach einmaligen Ansichten'
        toplist_downloads_slide.placeholders[1].text = f'Top {len(toplist_categories_info_list)} Categories per unique views'
        toplist_downloads_table =create_table(toplist_downloads_slide.placeholders[16], len(toplist_downloads_info_list)+1, 2)
        header_cell_1 = toplist_downloads_table.cell(0,0)
        header_cell_1.text = "Video"
        header_cell_2 = toplist_downloads_table.cell(0,1)
        header_cell_2.text = "Downloads"
        access_and_fill_table(toplist_downloads_table, toplist_downloads_info_list)
        make_text_bold_in_table(toplist_downloads_table, ['Video', 'Downloads'])
        toplist_downloads_slide.placeholders[13].text = f'unique views: from({user_info_dictionary['start_date']}) | til({user_info_dictionary['end_date']})'
        toplist_downloads_slide.placeholders[14].text = 'Source: Solidsport'
    elif user_info_dictionary['categories_exist'] == True:
        #toplist_downloads_slide.placeholders[0].text = 'Beliebteste Kategorien nach einmaligen Ansichten'
        toplist_downloads_slide.placeholders[1].text = f'Top {len(toplist_categories_info_list)} Categories per unique views'
        toplist_categories_table = create_table(toplist_downloads_slide.placeholders[16], len(toplist_categories_info_list)+1, 2)
        header_cell_1 = toplist_categories_table.cell(0,0)
        header_cell_1.text = "Categories"
        header_cell_2 = toplist_categories_table.cell(0,1)
        header_cell_2.text = "unique views"
        access_and_fill_table(toplist_categories_table, toplist_categories_info_list)
        make_text_bold_in_table(toplist_categories_table, ['Categorie', 'unique views'])
        toplist_downloads_slide.placeholders[13].text = f'unique views: from({user_info_dictionary['start_date']}) | till({user_info_dictionary['end_date']})'
        toplist_downloads_slide.placeholders[14].text = 'Source: Solidsport'

    # section_earnings slide
    seciton_earnings_slide.placeholders[14].text = 'Source: Solidsport'

    #overview_earnings slide
    if user_info_dictionary['multi_currrency'] == False and len(all_access_data) == 1 and len(settlements_data) == 1:
        overview_earnings_slide.placeholders[14].text = 'Source: Solidsport'

        earnings_table = create_revenue_table_1(overview_earnings_slide.placeholders[15],
                                            9, 3, all_access_data)
        fill_reavenue_table_1(earnings_table, user_info_dictionary, ppv_info_list,
                            billing_info, all_access_data, single_stream_purchase_count,
                            0)

        format_table(earnings_table, ['Categorie', 'Description', 'TOTAL', 
                                        'Total gross*', 'Net reavenue*'])
        analytics_table_2 = create_analytics_table_1(overview_earnings_slide.placeholders[18], 3, 1)
        header_cell_analyse_1 = analytics_table_2.cell(0,0)
        header_cell_analyse_1.text = "All Access passes / single streams"

        analyse_2_cell_10 = analytics_table_2.cell(1,0)
        analyse_2_cell_10.text = f"{total_all_access_passes} / {single_stream_purchase_count}  ≈ {calculations[2]}"
        analyse_2_cell_20 = analytics_table_2.cell(2,0)
        analyse_2_cell_20.text = f"roughly {calculations[2]} times all access passes per single streams were purchased"

        format_analytics_table(analytics_table_2, ["All Access Pässe / Einzelstreams"])

        overview_earnings_slide.placeholders[19].text = f'The avarage vat is: {calculations[3]*100}%'
    else:
        earnings_table = create_revenue_table_2(multi_currency_slide.placeholders[15],
                                                7, 9, all_access_data, 
                                                settlements_data)
        fill_reavenue_table_2(earnings_table, user_info_dictionary, ppv_info_list, 
                            billing_info, all_access_data, 
                            single_stream_purchase_count, multi_currency_info,
                            settlements_data, transaction_data)
        format_table(earnings_table, ['Categorie', 'Description', 'TOTAL', 
                                        'Total gross*', 'Net reavenue*', '∑'])

    # billing slide
    billing_slide.placeholders[14].text = 'Source: Solidsport'
    billing_slide.placeholders[19].text = billing_info['Pay Per View reimbursement up until']
    billing_slide.placeholders[20].text = billing_info['reference_number']
    billing_slide.placeholders[21].text = billing_info['amount']
    billing_slide.placeholders[22].text = billing_info['conditions_of_payment']

    #further_info slide

    # save Powerpoint
    folder_path_prs = os.path.join(os.path.expanduser("~"), "Desktop")
    prs.save(f'{folder_path_prs}/Analysereport {user_info_dictionary['event_name']}.pptx')

    print(f'PowerPoint-Präsentation: {user_info_dictionary['event_name']} wurde '
          'unter: {folder_path_prs} erfolgreich erstellt!')

    # delete folder
    erase_folder(user_info_dictionary)

# GUI functions
def create_presentation():
    """Creates a new presentation with the name "Analysereport [Event Name]"""
    if (event_name_entry.get() == '' or 
        start_date_entry.get() == '' or
        end_date_entry.get() == '' or
        page_url_entry.get() == '' or
        e_mail_entry.get() == '' or
        password_entry.get() == ''):
        error_message = ttk.Label(master=window,
                                  text="Bitte alle Eingabefelder ausfüllen!",
                                  font= "Calibri 20 bold",
                                  foreground='red')
        error_message.grid(row=5, column=0, pady=15)  
    else:
        main()
def retrieve_info():
    if (event_name_entry.get() == '' or 
        start_date_entry.get() == '' or
        end_date_entry.get() == '' or
        page_url_entry.get() == '' or
        e_mail_entry.get() == '' or
        password_entry.get() == ''):
        error_message = ttk.Label(master=window,
                                text="Bitte alle Eingabefelder ausfüllen!",
                                font= "Calibri 20 bold",
                                foreground='red')
        error_message.grid(row=5, column=0, pady=15)
    return (event_name_value.get(),
                page_url_value.get(),
                e_mail_value.get(),
                password_value.get(),
                start_date_value.get(),
                end_date_value.get(),
                revenue_split_value.get(),
                price_single_stream_value.get(),
                all_access_pass_exist_var.get(),
                multi_currency_var.get(),
                categories_exist_var.get(),
                downloads_exist_var.get(),
                powered_by_var.get())

# create a window
window = tk.Tk()
window.title('Analytics Automation')
window.geometry('450x800')

# titel
title_label = ttk.Label(master = window, 
                        text = 'Folgende Informationen müssen angegeben werden. \n' 
                        'Achte auf genaue Rechtschreibung!', 
                        font= 'Calibri 14 bold')
title_label.grid(column=0, row=0, pady= 10)

# Frame for entry fields for user information
userinformation_frame = ttk.Frame(master=window)
userinformation_frame.grid(column=0, row=1, pady=10)

# Entry fields for user information
event_name_label = ttk.Label(master=userinformation_frame,
                            text='Name of the event: ')
event_name_value = tk.StringVar()
event_name_entry = ttk.Entry(master=userinformation_frame,
                            textvariable= event_name_value)

page_url_label = ttk.Label(master=userinformation_frame,
                           text='Enter correct page URL here:')
page_url_value = tk.StringVar()
page_url_entry = ttk.Entry(master=userinformation_frame,
                           textvariable=page_url_value)

e_mail_label = ttk.Label(master=userinformation_frame,
                         text='Your Solidsport E-Mail: ')
e_mail_value = tk.StringVar()
e_mail_entry = ttk.Entry(master=userinformation_frame,
                         textvariable=e_mail_value)

password_label = ttk.Label(master=userinformation_frame,
                           text= 'Your Solidsport password: ')
password_value = tk.StringVar()
password_entry = ttk.Entry(master=userinformation_frame,
                           textvariable=password_value)

start_date_label = ttk.Label(master=userinformation_frame,
                             text= 'The starting date you want analytics to show:')
start_date_value = tk.StringVar()
start_date_entry = ttk.Entry(master=userinformation_frame,
                             textvariable=start_date_value)

end_date_label = ttk.Label(master=userinformation_frame,
                           text='The end date you want analytics to show:')
end_date_value = tk.StringVar()
end_date_entry = ttk.Entry(master=userinformation_frame,
                           textvariable=end_date_value)

revenue_split_label = ttk.Label(master=userinformation_frame,
                                text='Enter the customers revenue share '
                                '(Format: eg. 60):')
revenue_split_value = tk.StringVar()
reavenue_split_entry = ttk.Entry(master=userinformation_frame,
                                 textvariable=revenue_split_value)

price_single_stream_label =  ttk.Label(master=userinformation_frame,
                                       text= 'enter price single stream:')
price_single_stream_value = tk.StringVar()
price_single_stream_entry = ttk.Entry(master=userinformation_frame,
                                      textvariable=price_single_stream_value)

all_access_pass_exist_var = tk.BooleanVar(value=False)
all_access_pass_exist = ttk.Checkbutton(master=userinformation_frame, 
                                   text="do All Access Passess exit?",
                                   variable=all_access_pass_exist_var)
multi_currency_var = tk.BooleanVar(value=False)
multi_currency = ttk.Checkbutton(master=userinformation_frame,
                                 text= "Is this a multi currency channel?",
                                 variable=multi_currency_var)


# pack entry fields
event_name_label.pack()
event_name_entry.pack()
page_url_label.pack()
page_url_entry.pack()
e_mail_label.pack()
e_mail_entry.pack()
password_label.pack()
password_entry.pack()
start_date_label.pack()
start_date_entry.pack()
end_date_label.pack()
end_date_entry.pack()
revenue_split_label.pack()
reavenue_split_entry.pack()
price_single_stream_label.pack()
price_single_stream_entry.pack()
all_access_pass_exist.pack(pady=15)
multi_currency.pack(pady=0)

# Frame for Checkbuttons
checkbuttons_frame = ttk.Frame(master=window)
checkbuttons_frame.grid(column=0, row=2, pady=10)

# Checkbuttons
categories_exist_var = tk.BooleanVar(value=False)
categories_exist = ttk.Checkbutton(master=checkbuttons_frame, 
                                   text="Do you want top categories to show in the presentation?",
                                   variable=categories_exist_var)
downloads_exist_var = tk.BooleanVar(value=False)
downloads_exist = ttk.Checkbutton(master=checkbuttons_frame,
                                  text="Do you want the most downloaded videos to show in the presentation?",
                                  variable=downloads_exist_var)
powered_by_var = tk.BooleanVar(value=False)
powered_by = ttk.Checkbutton(master=checkbuttons_frame,
                             text="Is this a powered by website?",
                             variable=powered_by_var)

# pack Chekcbuttons
categories_exist.pack()
downloads_exist.pack()
powered_by.pack()

# retrieve info button
retrieve_button = ttk.Button(master=window,
                             text='Retrieve info',
                             command=retrieve_info)
retrieve_button.grid(column=0, row=3, pady=20)

# create presentation button
create_presentation_button = ttk.Button(master=window, 
                                        text="Create presentation",
                                        command=create_presentation)
create_presentation_button.grid(column=0, row=4)

print("Process finished --- %s seconds ---" % (time.time() - start_time))

# run window
window.mainloop()